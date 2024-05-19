from fastapi import FastAPI
from pptx import Presentation
import re
from langchain import PromptTemplate
from langchain_google_genai import ChatGoogleGenerativeAI

app = FastAPI()

def extract_text_from_pptx(file_path):
    try:
        presentation = Presentation(file_path)
        text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        print(f"Error extracting text from PowerPoint: {e}")
        return ""

@app.post("/process_presentation")
async def process_presentation(file_path: str):
    extracted_text = extract_text_from_pptx(file_path)
    cleaned_text = re.sub(r'\s+', ' ', extracted_text)

    generation_config = {
        "temperature": 0.3,
        "top_p": 0.8,
        "top_k": 15,
        "max_output_tokens": 2048,
    }

    llm = ChatGoogleGenerativeAI(
        model="gemini-pro",
        google_api_key="YOUR_GOOGLE_API_KEY",
        generation_config=generation_config
    )

    presentation = cleaned_text

    template = (
        "Tell me the topic which belongs to this presentation: {presentation} \n\n"
        "Also evaluate the presentation on a scale from 20% to 100% according to presentation \n\n"
        "Additionally, provide a summary of the presentation."
    )

    prompt = PromptTemplate(
        input_variables=["presentation"],
        template=template,
    )

    formatted_prompt = prompt.format(presentation=presentation)

    response = llm.invoke(formatted_prompt)

    response_content = re.sub(r'\*\*(.*?)\*\*', r'\1', response.content)

    return {"response": response_content}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="127.0.0.1", port=8000)
